"""Build .pptx files from a JSON spec (CLI: python -m koraku.artifacts.pptx_build)."""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any


def _load_spec(path: str | None, inline: str | None) -> dict[str, Any]:
    if path:
        return json.loads(Path(path).read_text(encoding="utf-8"))
    if inline:
        return json.loads(inline)
    return {}


from koraku.artifacts.pptx_layout import normalize_presentation_spec, slide_body_lines


def build_pptx(spec: dict[str, Any], out_path: str) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as e:
        raise SystemExit(
            "python-pptx is required. Install with: pip install 'koraku[artifacts]'"
        ) from e

    prs = Presentation()
    spec = normalize_presentation_spec(spec)
    title_slide = str(spec.get("title") or "").strip()
    subtitle = str(spec.get("subtitle") or "").strip()
    if title_slide:
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title_slide
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

    for slide_spec in spec.get("slides") or []:
        if not isinstance(slide_spec, dict):
            continue
        layout_idx = int(slide_spec.get("layout") or 1)
        layout_idx = min(max(layout_idx, 0), len(prs.slide_layouts) - 1)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        stitle = str(slide_spec.get("title") or "").strip()
        if stitle and slide.shapes.title:
            slide.shapes.title.text = stitle

        body = slide_body_lines(slide_spec)

        if body and len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, line in enumerate(body):
                text = str(line).strip()
                if not text:
                    continue
                if i == 0:
                    tf.text = text
                else:
                    p = tf.add_paragraph()
                    p.text = text
                    p.level = 0
                    p.font.size = Pt(18)

        notes = str(slide_spec.get("notes") or "").strip()
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    out = Path(out_path).expanduser().resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Build a .pptx from JSON spec")
    parser.add_argument("--out", required=True, help="Output .pptx path")
    parser.add_argument("--spec", help="Path to JSON spec file")
    parser.add_argument("--json", dest="inline_json", help="Inline JSON spec")
    args = parser.parse_args(argv)

    spec = _load_spec(args.spec, args.inline_json)
    path = build_pptx(spec, args.out)
    print(json.dumps({"ok": True, "path": path, "type": "pptx", "slides": len(spec.get("slides") or [])}))
    return 0


if __name__ == "__main__":
    sys.exit(main())
