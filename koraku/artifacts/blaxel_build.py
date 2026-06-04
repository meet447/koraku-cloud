"""Run artifact builders inside a Blaxel sandbox when host file tools are routed there."""
from __future__ import annotations

import base64
import json
import shlex
from typing import Any

from koraku.agent.blaxel_scope import get_active_blaxel_session_root
from koraku.tools.blaxel_dispatch import blaxel_bash_if_active, sandbox_abs_path


def _pip_install(package: str) -> str:
    return (
        f"(pip install -q {shlex.quote(package)} 2>/dev/null "
        f"|| pip3 install -q {shlex.quote(package)} 2>/dev/null "
        f"|| pip3 install -q --break-system-packages {shlex.quote(package)} 2>/dev/null "
        f"|| true)"
    )


def _blaxel_pptx_script(spec: dict[str, Any], abs_out: str) -> str:
    b64 = base64.b64encode(json.dumps(spec).encode()).decode()
    return f"""{_pip_install("python-pptx")}
python3 <<'PY'
import base64, json, os
from pptx import Presentation
from pptx.util import Pt

spec = json.loads(base64.b64decode({json.dumps(b64)}).decode())
out = {json.dumps(abs_out)}
os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
prs = Presentation()
title = str(spec.get("title") or "").strip()
subtitle = str(spec.get("subtitle") or "").strip()
if title:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
for slide_spec in spec.get("slides") or []:
    if not isinstance(slide_spec, dict):
        continue
    layout_idx = min(max(int(slide_spec.get("layout") or 1), 0), len(prs.slide_layouts) - 1)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    stitle = str(slide_spec.get("title") or "").strip()
    if stitle and slide.shapes.title:
        slide.shapes.title.text = stitle
    body = slide_spec.get("body") or slide_spec.get("bullets") or []
    if isinstance(body, str):
        body = [line.strip() for line in body.split("\\n") if line.strip()]
    if body and len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, line in enumerate(body):
            text = str(line).strip()
            if not text:
                continue
            if i == 0:
                tf.text = text
            else:
                p = tf.add_paragraph()
                p.text = text
                p.level = 0
                p.font.size = Pt(18)
prs.save(out)
print(json.dumps({{"ok": True, "path": out, "type": "pptx", "slides": len(prs.slides)}}))
PY"""


def _blaxel_docx_script(spec: dict[str, Any], abs_out: str) -> str:
    b64 = base64.b64encode(json.dumps(spec).encode()).decode()
    return f"""{_pip_install("python-docx")}
python3 <<'PY'
import base64, json, os
from docx import Document

spec = json.loads(base64.b64decode({json.dumps(b64)}).decode())
out = {json.dumps(abs_out)}
os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
doc = Document()
title = str(spec.get("title") or "").strip()
if title:
    doc.add_heading(title, level=0)
for section in spec.get("sections") or []:
    if not isinstance(section, dict):
        continue
    heading = str(section.get("heading") or section.get("title") or "").strip()
    if heading:
        doc.add_heading(heading, level=min(max(int(section.get("level") or 1), 1), 3))
    body = section.get("body") or section.get("text") or ""
    if isinstance(body, list):
        for item in body:
            if str(item).strip():
                doc.add_paragraph(str(item).strip(), style="List Bullet")
    elif str(body).strip():
        for para in str(body).split("\\n\\n"):
            if para.strip():
                doc.add_paragraph(para.strip())
    for row in section.get("bullets") or []:
        if str(row).strip():
            doc.add_paragraph(str(row).strip(), style="List Bullet")
doc.save(out)
print(json.dumps({{"ok": True, "path": out, "type": "docx"}}))
PY"""


def _blaxel_xlsx_script(spec: dict[str, Any], abs_out: str) -> str:
    b64 = base64.b64encode(json.dumps(spec).encode()).decode()
    return f"""{_pip_install("openpyxl")}
python3 <<'PY'
import base64, json, os
from openpyxl import Workbook

spec = json.loads(base64.b64decode({json.dumps(b64)}).decode())
out = {json.dumps(abs_out)}
os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
wb = Workbook()
sheets = spec.get("sheets") or [{{"name": spec.get("sheet_name") or "Sheet1", "headers": spec.get("headers") or [], "rows": spec.get("rows") or []}}]
first = True
for sheet_spec in sheets:
    if not isinstance(sheet_spec, dict):
        continue
    if first:
        ws = wb.active
        ws.title = str(sheet_spec.get("name") or "Sheet1")[:31]
        first = False
    else:
        ws = wb.create_sheet(title=str(sheet_spec.get("name") or "Sheet")[:31])
    headers = sheet_spec.get("headers") or []
    if headers:
        for ci, header in enumerate(headers, start=1):
            ws.cell(row=1, column=ci, value=str(header))
    start_row = 2 if headers else 1
    for ri, row in enumerate(sheet_spec.get("rows") or [], start=start_row):
        if not isinstance(row, list):
            ws.cell(row=ri, column=1, value=str(row))
            continue
        for ci, val in enumerate(row, start=1):
            ws.cell(row=ri, column=ci, value=val)
wb.save(out)
print(json.dumps({{"ok": True, "path": out, "type": "xlsx"}}))
PY"""


def _parse_build_json(result: str) -> dict[str, Any] | None:
    for line in reversed(result.splitlines()):
        line = line.strip()
        if not line.startswith("{"):
            continue
        try:
            payload = json.loads(line)
        except json.JSONDecodeError:
            continue
        if isinstance(payload, dict) and payload.get("ok"):
            return payload
    return None


async def blaxel_build_artifact(artifact_type: str, spec: dict[str, Any], output_rel: str) -> str:
    out_rel = (output_rel or "").strip().replace("\\", "/")
    if not out_rel:
        return "Error: output_path is required."

    session_root = (get_active_blaxel_session_root() or "").strip()
    if not session_root:
        return (
            "Error: chat session folder is not bound in the sandbox. "
            "Retry after the sandbox attaches (first file tool on this turn)."
        )

    abs_out = sandbox_abs_path(out_rel)

    if artifact_type == "presentation":
        cmd = _blaxel_pptx_script(spec, abs_out)
    elif artifact_type == "document":
        cmd = _blaxel_docx_script(spec, abs_out)
    elif artifact_type == "spreadsheet":
        cmd = _blaxel_xlsx_script(spec, abs_out)
    else:
        return f"Error: unsupported Blaxel artifact type {artifact_type!r}"

    result = await blaxel_bash_if_active(cmd, timeout=120)
    if result is None:
        return "Error: Blaxel sandbox is not active."
    if result.strip().startswith("Error"):
        return result.strip()

    payload = _parse_build_json(result)
    if payload is None:
        return f"Error: artifact build did not produce a success payload.\n{result[:2000]}"

    verify = await blaxel_bash_if_active(
        f"test -f {shlex.quote(abs_out)} && echo __KORAKU_OK__",
        timeout=30,
    )
    if verify is None or "__KORAKU_OK__" not in verify:
        return (
            f"Error: build reported success but file is missing in session folder.\n"
            f"Expected: {out_rel}\nSession root: {session_root}"
        )

    payload["path"] = out_rel
    payload["session_root"] = session_root
    return json.dumps(payload)


async def blaxel_merge_pdfs(inputs: list[str], output_rel: str) -> str:
    out_rel = (output_rel or "").strip().replace("\\", "/")
    session_root = (get_active_blaxel_session_root() or "").strip()
    if not session_root:
        return "Error: chat session folder is not bound in the sandbox."
    abs_out = sandbox_abs_path(out_rel)
    abs_inputs = [sandbox_abs_path(str(p).replace("\\", "/")) for p in inputs]
    cmd = f"""{_pip_install("pypdf")}
python3 <<'PY'
import json
from pypdf import PdfWriter
writer = PdfWriter()
for path in {json.dumps(abs_inputs)}:
    writer.append(path)
out = {json.dumps(abs_out)}
import os
os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
with open(out, "wb") as f:
    writer.write(f)
print(json.dumps({{"ok": True, "path": out, "type": "pdf"}}))
PY"""
    result = await blaxel_bash_if_active(cmd, timeout=120)
    if result is None:
        return "Error: Blaxel sandbox is not active."
    payload = _parse_build_json(result or "")
    if payload is None:
        return result.strip() or "Error: empty response from Blaxel PDF merge."
    verify = await blaxel_bash_if_active(
        f"test -f {shlex.quote(abs_out)} && echo __KORAKU_OK__",
        timeout=30,
    )
    if verify is None or "__KORAKU_OK__" not in verify:
        return f"Error: merged PDF missing at {out_rel}"
    payload["path"] = out_rel
    return json.dumps(payload)
